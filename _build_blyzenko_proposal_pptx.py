#!/usr/bin/env python3
# Генерує Blyzenko-Bolt-Food-Commercial-Proposal.pptx зі змісту blyzenko-commercial-proposal.html
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

GREEN = RGBColor(0x34, 0xD1, 0x86)
GREEN_D = RGBColor(0x27, 0xB9, 0x75)
DARK = RGBColor(0x2F, 0x31, 0x3F)
GREY = RGBColor(0x6B, 0x70, 0x80)
LIGHT = RGBColor(0xEE, 0xF6, 0xF1)
BARLT = RGBColor(0xBF, 0xE9, 0xD3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DSUB = RGBColor(0xC9, 0xCE, 0xD6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = 13.333


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); r.shadow.inherit = False
    return s


def tbox(s, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, space=0):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for ln in lines:
        text, size, color, bold = ln
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = align
        if space and not first:
            p.space_before = Pt(space)
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color; run.font.name = 'Arial'
        first = False
    return tb


def rrect(s, l, t, w, h, fill, rad=0.08, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = rad
    except Exception:
        pass
    return sp


def rect(s, l, t, w, h, fill):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def kicker(s, kick, title, dark=False):
    rect(s, 0.7, 0.62, 0.55, 0.1, GREEN)
    tbox(s, 0.7, 0.78, 12, 0.35, [(kick.upper(), 12, GREEN, True)])
    tbox(s, 0.7, 1.08, 12, 1.0, [(title, 27, WHITE if dark else DARK, True)])


def card(s, l, t, w, h, title, body):
    rrect(s, l, t, w, h, LIGHT, rad=0.05)
    rect(s, l, t + 0.04, 0.07, h - 0.08, GREEN)
    tbox(s, l + 0.28, t + 0.2, w - 0.5, h - 0.32,
         [(title, 15, DARK, True), (body, 10.5, GREY, False)], space=6)


def tile(s, l, t, w, h, big, lab, bigcolor=GREEN):
    rrect(s, l, t, w, h, LIGHT, rad=0.08)
    tbox(s, l + 0.1, t + 0.28, w - 0.2, 0.85, [(big, 23, bigcolor, True)], align=PP_ALIGN.CENTER)
    tbox(s, l + 0.12, t + h - 0.95, w - 0.24, 0.85, [(lab, 10.5, GREY, False)], align=PP_ALIGN.CENTER)


def row(s, l, t, w, left, right=None, h=0.52):
    rrect(s, l, t, w, h, LIGHT, rad=0.2)
    if right is None:
        tbox(s, l + 0.2, t, w - 0.4, h, [(left, 11.5, DARK, False)], anchor=MSO_ANCHOR.MIDDLE)
    else:
        tbox(s, l + 0.2, t, w * 0.62, h, [(left, 11.5, DARK, False)], anchor=MSO_ANCHOR.MIDDLE)
        tbox(s, l + w * 0.5, t, w * 0.45, h, [(right, 11.5, GREEN, True)], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)


def boldrow(s, l, t, w, b, rest, h=0.62):
    rrect(s, l, t, w, h, LIGHT, rad=0.16)
    tb = s.shapes.add_textbox(Inches(l + 0.2), Inches(t), Inches(w - 0.4), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = b + ' '; r1.font.size = Pt(11); r1.font.bold = True; r1.font.color.rgb = DARK; r1.font.name = 'Arial'
    r2 = p.add_run(); r2.text = rest; r2.font.size = Pt(11); r2.font.bold = False; r2.font.color.rgb = GREY; r2.font.name = 'Arial'


def step(s, l, t, w, n, text):
    b = rrect(s, l, t + 0.02, 0.34, 0.34, GREEN, rad=0.25)
    tbox(s, l, t + 0.02, 0.34, 0.34, [(str(n), 12, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tbox(s, l + 0.5, t, w - 0.5, 0.45, [(text, 11.5, DARK, False)], anchor=MSO_ANCHOR.MIDDLE)


def chip(s, l, t, w, text):
    rrect(s, l, t, w, 0.5, LIGHT, rad=0.25)
    tbox(s, l + 0.18, t, w - 0.36, 0.5, [(text, 10.5, DARK, False)], anchor=MSO_ANCHOR.MIDDLE)


def bigmsg(s, l, t, w, text, h=0.7):
    rrect(s, l, t, w, h, GREEN, rad=0.12)
    tbox(s, l + 0.25, t, w - 0.5, h, [(text, 12, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def note(s, text, t=6.55):
    tbox(s, 0.7, t, 11.93, 0.7, [(text, 11, GREY, False)])


def foot(s, text):
    tbox(s, 0.7, 7.05, 11.93, 0.35, [(text, 8.5, GREY, False)])


M = 0.7
CW = 11.93

# ===== 1 title =====
s = slide(DARK)
tbox(s, 0, 2.5, W, 0.4, [("КОМЕРЦІЙНА ПРОПОЗИЦІЯ", 13, GREEN, True)], align=PP_ALIGN.CENTER)
tbox(s, 0, 2.95, W, 1.0, [("Близенько × Bolt Food", 40, WHITE, True)], align=PP_ALIGN.CENTER)
tbox(s, 1.5, 4.05, W - 3, 0.8, [("Якірний рітейл регіону — другий канал, що додає оборот без канібалізації", 16, DSUB, False)], align=PP_ALIGN.CENTER)
tbox(s, 0, 5.0, W, 0.4, [("Україна · Червень 2026", 13, RGBColor(0x9A, 0xA0, 0xAC), False)], align=PP_ALIGN.CENTER)

# ===== 2 ecosystem =====
s = slide()
kicker(s, "Meet Bolt", "Вся екосистема в одному застосунку")
eco = [("Таксі", "40+ країн", "500+ міст"), ("Доставка", "15+ країн", "200+ міст"), ("Оренда / самокати", "20+ країн", "200+ міст")]
ew = (CW - 2 * 0.3) / 3
for i, (hd, a, bb) in enumerate(eco):
    l = M + i * (ew + 0.3); t = 2.4; h = 2.3
    rrect(s, l, t, ew, h, LIGHT, rad=0.07)
    rrect(s, l, t, ew, 0.65, GREEN, rad=0.07)
    tbox(s, l, t, ew, 0.65, [(hd, 15, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tbox(s, l, t + 0.65, ew, h - 0.65, [(a, 18, DARK, True), (bb, 18, DARK, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space=4)
note(s, "Крос-вертикальний ефект: користувачів таксі та Food ми конвертуємо у покупців Близенько — нові цифрові клієнти поза вашою офлайн-базою.", t=5.2)

# ===== 3 bolt numbers =====
s = slide()
kicker(s, "Bolt today", "Bolt — у цифрах")
nums = [("50+", "країн"), ("850+", "міст"), ("4 млн+", "замовлень/день"), ("€12 млрд+", "GMV run-rate"), ("30%+", "GMV CAGR 21–25")]
tw = (CW - 4 * 0.18) / 5
for i, (big, lab) in enumerate(nums):
    tile(s, M + i * (tw + 0.18), 2.5, tw, 2.0, big, lab)
note(s, "Один застосунок: доставка їжі та товарів, таксі, оренда — мільйони активних клієнтів в Україні.", t=5.0)

# ===== 4 stores growth =====
s = slide()
kicker(s, "Сегмент stores у Bolt Food", "Напрям магазинів зростає шість місяців поспіль")
cbl, cbt, cbw, cbh = M, 2.4, 7.0, 3.7
rrect(s, cbl, cbt, cbw, cbh, LIGHT, rad=0.05)
tbox(s, cbl + 0.3, cbt + 0.25, cbw - 0.6, 0.4, [("▲ стабільне зростання 6 міс поспіль", 11, GREEN, True)])
bars = [("–5 міс", 26), ("–4", 39), ("–3", 53), ("–2", 69), ("–1", 84), ("Зараз", 100)]
base = cbt + cbh - 0.55
maxh = 2.4
bw = (cbw - 0.6) / len(bars) * 0.62
gap = (cbw - 0.6) / len(bars)
for i, (cap, pct) in enumerate(bars):
    bh = maxh * pct / 100
    bl = cbl + 0.35 + i * gap + (gap - bw) / 2
    rrect(s, bl, base - bh, bw, bh, BARLT, rad=0.04)
    rect(s, bl, base - bh, bw, 0.07, GREEN)
    tbox(s, bl - 0.15, base + 0.05, bw + 0.3, 0.3, [(cap, 9, GREY, False)], align=PP_ALIGN.CENTER)
scl = cbl + cbw + 0.35; scw = CW - cbw - 0.35
sc = [("stores", "найшвидше зростаючий сегмент Bolt Food"), ("×4", "зростання активних локацій за 6 міс"), ("↗ частка", "напрям збільшує вагу в сегменті")]
for i, (v, t) in enumerate(sc):
    ct = 2.4 + i * 1.28
    rrect(s, scl, ct, scw, 1.12, LIGHT, rad=0.08)
    tbox(s, scl + 0.25, ct + 0.16, scw - 0.4, 0.5, [(v, 18, GREEN, True)])
    tbox(s, scl + 0.25, ct + 0.6, scw - 0.4, 0.45, [(t, 9.5, GREY, False)])
foot(s, "Активні локації stores (стовпчики) · частка stores у сегменті. Дані ілюстративні, без точних значень.")

# ===== 5 why bolt =====
s = slide()
kicker(s, "Що ми даємо як компанія", "Чому Bolt для Близенько")
why = [
    ("100% інкрементальний канал", "Близенько = 0 на Bolt Food. Кожне замовлення — поверх наявного каналу, без канібалізації + диверсифікація залежності від однієї платформи."),
    ("Профіль під q-commerce", "Продукти біля дому — висока частота й повторні покупки. Ідеально під Bolt+ і безкоштовну доставку."),
    ("Дірект-інтеграція POS", "Ваші реальні ціни, синхронізація стоків, контроль асортименту й промо. Capex на інтеграцію — на нас."),
    ("Регіональна сила на Заході", "Розширюємось із Рукавичкою → щільніша мережа кур'єрів і більше попиту в домашньому регіоні Близенько."),
]
cw = (CW - 0.3) / 2
for i, (ti, bo) in enumerate(why):
    l = M + (i % 2) * (cw + 0.3); t = 2.4 + (i // 2) * 2.05
    card(s, l, t, cw, 1.9, ti, bo)

# ===== 6 promo tools =====
s = slide()
kicker(s, "Промо · як розганяємо замовлення", "Промо-інструменти Bolt Food")
promo = [("1", "Знижка на доставку", "Частково або 0 ₴. №1 інструмент для залучення нових клієнтів."),
         ("2", "Знижки на каталог", "Окремі позиції або весь асортимент — через API або від Bolt."),
         ("3", "Когорти", "Таргет на Нових / Втрачених / Активних — різний розмір знижки."),
         ("4", "Банерна підтримка", "Банер у вкладці «Магазини» в періоди активних акцій.")]
pw = (CW - 3 * 0.18) / 4
for i, (n, h4, p) in enumerate(promo):
    l = M + i * (pw + 0.18); t = 2.4
    rrect(s, l, t, pw, 2.1, LIGHT, rad=0.06)
    rrect(s, l + 0.22, t + 0.22, 0.4, 0.4, GREEN, rad=0.2)
    tbox(s, l + 0.22, t + 0.22, 0.4, 0.4, [(n, 12, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tbox(s, l + 0.22, t + 0.78, pw - 0.44, 0.5, [(h4, 12, DARK, True)])
    tbox(s, l + 0.22, t + 1.25, pw - 0.44, 0.8, [(p, 9, GREY, False)])
pb = [("Рекомендований старт", "Безкоштовна доставка для нових клієнтів + MOV трохи вищий за середній чек.", GREEN),
      ("Знижки через API", "Передаєте фінальні ціни — без +20% ПДВ на знижку.", DARK)]
bw = (CW - 0.3) / 2
for i, (h4, p, col) in enumerate(pb):
    l = M + i * (bw + 0.3); t = 4.8
    rrect(s, l, t, bw, 1.2, LIGHT, rad=0.06)
    rect(s, l, t + 0.04, 0.07, 1.12, col)
    tbox(s, l + 0.28, t + 0.2, bw - 0.5, 0.4, [(h4, 12, DARK, True)])
    tbox(s, l + 0.28, t + 0.62, bw - 0.5, 0.5, [(p, 9.5, GREY, False)])

# ===== 7 launch broad campaign =====
s = slide()
kicker(s, "Підтримка запуску", "Запуск підтримаємо широкою кампанією")
tbox(s, M, 1.95, CW, 0.4, [("Преміальні розміщення всередині застосунку — там, де клієнт ухвалює рішення про замовлення.", 13, GREY, False)])
cw = (CW - 0.4) / 2
cols = [("In-app модалки, банери та соцмережі", ["In-app modal — повноекранне промо при вході", "In-app банер у стрічці магазинів", "Social media — пости та сторіз Bolt Food"]),
        ("Промо-банери та категорії на головній", ["Банер «Безкоштовна доставка» на головному екрані", "Спец-категорія / окрема полиця Близенько", "Видимість у періоди акцій і сезонів"])]
for i, (ct, items) in enumerate(cols):
    l = M + i * (cw + 0.4); t = 2.7
    tbox(s, l, t, cw, 0.4, [(ct, 14, DARK, True)])
    for j, it in enumerate(items):
        row(s, l, t + 0.55 + j * 0.62, cw, it)
note(s, "Комбінація каналів Bolt Food + Bolt Rides дає охоплення, недоступне поодинці — до 10× на старті.", t=6.0)

# ===== 8 in-store branding =====
s = slide()
kicker(s, "In-store branding", "Підтримка офлайн-брендингу")
br = [("Брендинг у точках —", "Bolt підтримує оформлення та видимість бренду на місцях продажу."),
      ("Звʼязок онлайн ↔ офлайн —", "офлайн-аудиторія Близенько дізнається про вашу присутність на Bolt Food."),
      ("Спільні матеріали —", "банери, наліпки, POS-матеріали для промо-періодів.")]
for i, (b, rest) in enumerate(br):
    boldrow(s, M, 2.6 + i * 0.78, CW, b, rest)
note(s, "Онлайн і офлайн працюють разом: магазини Близенько стають точками входу нових клієнтів у застосунок.", t=5.6)

# ===== 9 our offer =====
s = slide()
kicker(s, "Наша пропозиція", "Комісія та безкоштовна доставка")
cyw = (5.5 - 0.3) / 2
for i, (yr, pct, nt) in enumerate([("2026", "6%", "комісія — стартовий рік"), ("2027", "8%", "комісія — другий рік")]):
    l = M + i * (cyw + 0.3); t = 2.3
    rrect(s, l, t, cyw, 1.5, LIGHT, rad=0.08)
    tbox(s, l, t + 0.12, cyw, 0.35, [(yr, 12, GREY, True)], align=PP_ALIGN.CENTER)
    tbox(s, l, t + 0.42, cyw, 0.7, [(pct, 30, GREEN, True)], align=PP_ALIGN.CENTER)
    tbox(s, l, t + 1.12, cyw, 0.35, [(nt, 9.5, GREY, False)], align=PP_ALIGN.CENTER)
tbox(s, M, 3.95, CW, 0.4, [("Безкоштовна доставка — спільне фінансування", 14, DARK, True)])
rows_tbl = [("Період", "Хто покриває доставку", "Логіка"),
            ("Перші 6 міс (старт)", "Bolt 100%", "якірний рітейл регіону — заводимо клієнтів за наш кошт"),
            ("Рік 2", "50% Bolt / 50% Близенько", "спільна інвестиція в утриманий попит"),
            ("Рік 3", "75% Близенько / 25% Bolt", "канал уже зрілий, Bolt лишається в долі"),
            ("Додатково", "OOH-підтримка", "зовнішня реклама на старті понад digital")]
gt = 4.45; gl = M; gw = CW; rh = 0.5
colw = [gw * 0.22, gw * 0.30, gw * 0.48]
for ri, rdata in enumerate(rows_tbl):
    ry = gt + ri * rh
    cx = gl
    for ci, cell in enumerate(rdata):
        if ri == 0:
            rect(s, cx, ry, colw[ci], rh, DARK)
            tbox(s, cx + 0.15, ry, colw[ci] - 0.3, rh, [(cell, 10, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
        else:
            bg = LIGHT if ri % 2 == 0 else WHITE
            rect(s, cx, ry, colw[ci], rh, bg)
            col = GREEN if ci == 1 else (DARK if ci == 0 else GREY)
            bold = ci <= 1
            tbox(s, cx + 0.15, ry, colw[ci] - 0.3, rh, [(cell, 9.5, col, bold)], anchor=MSO_ANCHOR.MIDDLE)
        cx += colw[ci]
foot(s, "Безкоштовна доставка 3–6 міс — обовʼязковий елемент старту для якірного рітейлу в регіоні.")

# ===== 10 Bolt+ =====
s = slide()
kicker(s, "Bolt+", "Преміум-аудиторія, що замовляє частіше")
bp = [("2.1–2.2×", "частіше замовляють підписники Bolt+", GREEN), ("+12–17%", "замовлень/міс після оформлення підписки", GREEN), ("3 сервіси", "таксі · Bolt Food · самокати в одній підписці", DARK)]
tw = (CW - 2 * 0.18) / 3
for i, (big, lab, c) in enumerate(bp):
    tile(s, M + i * (tw + 0.18), 2.3, tw, 1.7, big, lab, bigcolor=c)
rowsbp = [("Рік 1 — безкоштовно:", "Bolt+ підключено без жодної надбавки до комісії."),
          ("Рік 2 — опційно:", "за бажанням +1–2% лише на замовлення підписників; решта — за базовою ставкою."),
          ("Повна свобода:", "ви вільні відмовитись і не використовувати Bolt+ — це опція, а не зобовʼязання.")]
for i, (b, rest) in enumerate(rowsbp):
    boldrow(s, M, 4.25 + i * 0.62, CW, b, rest, h=0.54)
bigmsg(s, M, 6.25, CW, "Принцип: кожна 1 ₴ додаткової комісії за Bolt+ → ~7–8 ₴ додаткового обороту.")

# ===== 11 GMV potential =====
s = slide()
kicker(s, "Товарообіг · потенціал", "GMV-потенціал Близенько на Bolt")
chips = ["Старт: 60% покриття мережі → зростання до 100%", "середній чек 650 ₴", "орієнтир 2–3 замовлення/день на магазин"]
cx = M
for c in chips:
    cwid = 0.45 + len(c) * 0.078
    chip(s, cx, 2.2, cwid, c)
    cx += cwid + 0.2
scen = [("Базовий", "2 зам/день на магазин", "5,77 млн ₴", "≈ 69 млн ₴ / рік", False),
        ("Оптимальний", "3 зам/день на магазин", "8,66 млн ₴", "≈ 104 млн ₴ / рік", True)]
sw = (CW - 0.3) / 2
for i, (nm, sh, gmv, ordd, hot) in enumerate(scen):
    l = M + i * (sw + 0.3); t = 3.0; h = 2.5
    rrect(s, l, t, sw, h, GREEN if hot else LIGHT, rad=0.1)
    tbox(s, l, t + 0.3, sw, 0.4, [(nm, 16, WHITE if hot else DARK, True)], align=PP_ALIGN.CENTER)
    tbox(s, l, t + 0.75, sw, 0.35, [(sh, 11, RGBColor(0xEA, 0xFF, 0xF4) if hot else GREY, False)], align=PP_ALIGN.CENTER)
    tbox(s, l, t + 1.15, sw, 0.6, [(gmv, 26, WHITE if hot else GREEN, True)], align=PP_ALIGN.CENTER)
    tbox(s, l, t + 1.75, sw, 0.35, [("GMV / місяць", 11, RGBColor(0xEA, 0xFF, 0xF4) if hot else GREY, False)], align=PP_ALIGN.CENTER)
    tbox(s, l, t + 2.05, sw, 0.35, [(ordd, 12, WHITE if hot else DARK, True)], align=PP_ALIGN.CENTER)
note(s, "Прогноз — на старті при 60% покриття мережі; з розширенням до 100% потенціал зростає пропорційно. Повністю інкрементальний обсяг — новий канал поза наявним попитом.", t=5.85)
foot(s, "Припущення: середній чек 650 ₴, 30 днів/міс. Параметри перераховуються під фактичні дані Близенько.")

# ===== 12 next steps =====
s = slide()
kicker(s, "Наступні кроки", "Як швидко стартуємо")
cw = (CW - 0.4) / 2
tbox(s, M, 2.3, cw, 0.4, [("Що пропонуємо коротко", 14, DARK, True)])
offer = [("Комісія", "6% (2026) → 8% (2027)"), ("Доставка", "Bolt 100% перші 6 міс"), ("Далі", "50/50 рік 2 · 75/25 рік 3"), ("Bolt+", "рік 1 безкоштовно, далі опційно")]
for i, (a, b) in enumerate(offer):
    row(s, M, 2.85 + i * 0.62, cw, a, b)
l2 = M + cw + 0.4
tbox(s, l2, 2.3, cw, 0.4, [("Кроки", 14, DARK, True)])
steps = ["Підтвердити інтерес і scope пілоту", "Узгодити умови комісії та доставки", "Дірект-інтеграція каталогу та цін (capex Bolt)", "Запуск у регіоні + спільний маркетинг", "Огляд результатів через 4 тижні"]
for i, st in enumerate(steps):
    step(s, l2, 2.85 + i * 0.6, cw, i + 1, st)
bigmsg(s, M, 6.2, CW, "Готові запустити Близенько на Bolt Food — інкрементальний оборот, безкоштовна доставка на старті та Bolt+ безкоштовно в перший рік.")

import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Blyzenko-Bolt-Food-Commercial-Proposal.pptx")
prs.save(out)
print("saved", out, len(prs.slides.__iter__.__self__._sldIdLst))
